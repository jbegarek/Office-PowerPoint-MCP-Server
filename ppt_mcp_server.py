#!/usr/bin/env python
"""
MCP Server for PowerPoint manipulation using python-pptx.
Consolidated version with 20 tools organized into multiple modules.
"""
import os
import argparse
from pathlib import Path
from typing import Dict, Any

from mcp.server.fastmcp import FastMCP
from starlette.requests import Request
from starlette.responses import FileResponse, JSONResponse
from utils.http_auth import (
    APIKeyMiddleware,
    build_download_url,
    get_api_key,
    get_api_key_header_name,
)

# import utils  # Currently unused
from tools import (
    register_presentation_tools,
    register_content_tools,
    register_structural_tools,
    register_professional_tools,
    register_template_tools,
    register_hyperlink_tools,
    register_chart_tools,
    register_connector_tools,
    register_master_tools,
    register_transition_tools
)

# Output directory for saved presentations
OUTPUT_DIR = Path(os.environ.get("PPT_OUTPUT_DIR", "./output")).resolve()
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# Initialize the FastMCP server
app = FastMCP(
    name="ppt-mcp-server",
    instructions="""
You are connected to a PowerPoint presentation server.

WORKFLOW
1. Call create_presentation to start a new presentation (or create_presentation_from_templates for a full deck at once).
2. Add slides with add_slide or create_slide_from_template.
3. Add content to each slide: manage_text for text boxes, populate_placeholder / add_bullet_points for layout placeholders, add_table, add_chart, add_shape.
4. Optionally style with apply_professional_design or optimize_slide_text.
5. Call save_presentation as the final step — return the download_url from its response to the user.

FILE UPLOADS AND IMAGES — NOT SUPPORTED
- This server cannot receive any uploaded files. There is NO image tool available.
- Do NOT ask the user to upload images, templates, or any other files.
- Do NOT suggest adding images to slides — it is not possible.
- If the user asks to add an image or upload a file, tell them clearly: "This server cannot accept file uploads or add images."

CONVENTIONS
- All slide, shape, row, and column indices are 0-based (first = 0).
- Positions (left, top, width, height) are in inches.
- Colors are RGB lists: [R, G, B] where each value is 0–255. Example: [255, 0, 0] for red.
- presentation_id can almost always be omitted — the server remembers the current presentation.
- Available color_scheme values: "modern_blue", "corporate_gray", "elegant_green", "warm_red".

SAVING
- save_presentation must be called explicitly to write to disk and get a download link.
- Always call save_presentation last and return the download_url to the user.

TEXT TOOLS — WHEN TO USE WHICH
- populate_placeholder: fill a title/content/subtitle slot that exists in the slide's layout.
- add_bullet_points: fill a placeholder with a list of bullet items.
- manage_text (operation="add"): add a free-floating text box anywhere on the slide — use when you need text outside of layout placeholders.
- manage_text (operation="format"): change font, size, color, alignment on an existing shape.

TEMPLATE TOOLS — WHEN TO USE WHICH
- list_slide_templates: see all available built-in layout templates and their IDs.
- create_slide_from_template: PREFERRED — creates a new slide using a built-in template by ID.
- apply_slide_template: apply a built-in template to an already-existing slide.
- create_presentation_from_templates: create an entire presentation from a sequence of template IDs in one call.

TOOL CATEGORIES
- Presentation lifecycle:  create_presentation, save_presentation, list_presentations, switch_presentation, get_presentation_info, set_core_properties
- Add slides:              add_slide, create_slide_from_template, create_presentation_from_templates
- Text content:            manage_text, populate_placeholder, add_bullet_points
- Structural elements:     add_table, add_chart, add_shape, add_connector, format_table_cell, update_chart_data
- Design & styling:        apply_professional_design, optimize_slide_text
- Templates:               list_slide_templates, get_template_info, apply_slide_template, create_slide_from_template, create_presentation_from_templates
- Hyperlinks:              manage_hyperlinks
- Inspection:              get_slide_info, extract_slide_text, extract_presentation_text, manage_slide_masters
""",
)


# Global state to store presentations in memory
presentations = {}
current_presentation_id = None

# Transport mode (set in main() before app.run())
_transport_mode = "stdio"


def get_transport_mode() -> str:
    """Get the current transport mode."""
    return _transport_mode


def get_base_url() -> str:
    """Build the base URL from host/port settings.
    Uses RAILWAY_PUBLIC_DOMAIN when deployed on Railway."""
    public_domain = os.environ.get("RAILWAY_PUBLIC_DOMAIN")
    if public_domain:
        return f"https://{public_domain}"
    host = getattr(app.settings, "host", "localhost")
    port = getattr(app.settings, "port", 8000)
    display_host = "localhost" if host == "0.0.0.0" else host
    return f"http://{display_host}:{port}"


def get_download_url(filename: str) -> str:
    """Return a full download URL for a saved file."""
    safe_filename = Path(filename).name
    configured_base_url = os.environ.get("DOC_DOWNLOAD_BASE_URL", os.environ.get("MCP_DOWNLOAD_BASE_URL"))
    if configured_base_url:
        return build_download_url(configured_base_url, safe_filename)
    return build_download_url(f"{get_base_url()}/files", safe_filename)


def _build_streamable_http_app():
    """Build a streamable-http ASGI app across FastMCP versions."""
    if hasattr(app, "http_app"):
        try:
            return app.http_app(transport="streamable-http")
        except TypeError:
            return app.http_app()

    if hasattr(app, "streamable_http_app"):
        try:
            return app.streamable_http_app()
        except TypeError:
            return app.streamable_http_app(path="/mcp")

    raise RuntimeError("FastMCP version does not expose streamable HTTP app builder.")


def _build_sse_app():
    """Build an SSE ASGI app across FastMCP versions."""
    if hasattr(app, "http_app"):
        try:
            return app.http_app(transport="sse")
        except TypeError:
            pass

    if hasattr(app, "sse_app"):
        sse_path = getattr(getattr(app, "settings", None), "sse_path", "/sse")
        try:
            return app.sse_app(path=sse_path)
        except TypeError:
            try:
                return app.sse_app()
            except TypeError:
                pass

    raise RuntimeError("FastMCP version does not expose SSE app builder.")

# Template configuration
def get_template_search_directories():
    """
    Get list of directories to search for templates.
    Uses environment variable PPT_TEMPLATE_PATH if set, otherwise uses default directories.
    
    Returns:
        List of directories to search for templates
    """
    template_env_path = os.environ.get('PPT_TEMPLATE_PATH')
    
    if template_env_path:
        # If environment variable is set, use it as the primary template directory
        # Support multiple paths separated by colon (Unix) or semicolon (Windows)
        import platform
        separator = ';' if platform.system() == "Windows" else ':'
        env_dirs = [path.strip() for path in template_env_path.split(separator) if path.strip()]
        
        # Verify that the directories exist
        valid_env_dirs = []
        for dir_path in env_dirs:
            expanded_path = os.path.expanduser(dir_path)
            if os.path.exists(expanded_path) and os.path.isdir(expanded_path):
                valid_env_dirs.append(expanded_path)
        
        if valid_env_dirs:
            # Add default fallback directories
            return valid_env_dirs + ['.', './templates', './assets', './resources']
        else:
            print(f"Warning: PPT_TEMPLATE_PATH directories not found: {template_env_path}")
    
    # Default search directories when no environment variable or invalid paths
    return ['.', './templates', './assets', './resources']

# ---- Helper Functions ----

def get_current_presentation():
    """Get the current presentation object or raise an error if none is loaded."""
    if current_presentation_id is None or current_presentation_id not in presentations:
        raise ValueError("No presentation is currently loaded. Please create or open a presentation first.")
    return presentations[current_presentation_id]

def get_current_presentation_id():
    """Get the current presentation ID."""
    return current_presentation_id

def set_current_presentation_id(pres_id):
    """Set the current presentation ID."""
    global current_presentation_id
    current_presentation_id = pres_id

def validate_parameters(params):
    """
    Validate parameters against constraints.
    
    Args:
        params: Dictionary of parameter name: (value, constraints) pairs
        
    Returns:
        (True, None) if all valid, or (False, error_message) if invalid
    """
    for param_name, (value, constraints) in params.items():
        for constraint_func, error_msg in constraints:
            if not constraint_func(value):
                return False, f"Parameter '{param_name}': {error_msg}"
    return True, None

def is_positive(value):
    """Check if a value is positive."""
    return value > 0

def is_non_negative(value):
    """Check if a value is non-negative."""
    return value >= 0

def is_in_range(min_val, max_val):
    """Create a function that checks if a value is in a range."""
    return lambda x: min_val <= x <= max_val

def is_in_list(valid_list):
    """Create a function that checks if a value is in a list."""
    return lambda x: x in valid_list

def is_valid_rgb(color_list):
    """Check if a color list is a valid RGB tuple."""
    if not isinstance(color_list, list) or len(color_list) != 3:
        return False
    return all(isinstance(c, int) and 0 <= c <= 255 for c in color_list)

def add_shape_direct(slide, shape_type: str, left: float, top: float, width: float, height: float) -> Any:
    """
    Add an auto shape to a slide using direct integer values instead of enum objects.
    
    This implementation provides a reliable alternative that bypasses potential 
    enum-related issues in the python-pptx library.
    
    Args:
        slide: The slide object
        shape_type: Shape type string (e.g., 'rectangle', 'oval', 'triangle')
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        The created shape
    """
    from pptx.util import Inches
    
    # Direct mapping of shape types to their integer values
    # Values from MSO_AUTO_SHAPE_TYPE enum: https://github.com/scanny/python-pptx/blob/master/src/pptx/enum/shapes.py
    shape_type_map = {
        'rectangle': 1,              # RECTANGLE
        'rounded_rectangle': 5,      # ROUNDED_RECTANGLE
        'oval': 9,                   # OVAL
        'diamond': 4,                # DIAMOND
        'triangle': 7,               # ISOSCELES_TRIANGLE
        'right_triangle': 8,         # RIGHT_TRIANGLE
        'pentagon': 51,              # PENTAGON
        'hexagon': 10,               # HEXAGON
        'heptagon': 145,             # HEPTAGON
        'octagon': 6,                # OCTAGON
        'star': 92,                  # STAR_5_POINT
        'arrow': 33,                 # RIGHT_ARROW
        'cloud': 179,                # CLOUD
        'heart': 21,                 # HEART
        'lightning_bolt': 22,        # LIGHTNING_BOLT
        'sun': 23,                   # SUN
        'moon': 24,                  # MOON
        'smiley_face': 17,           # SMILEY_FACE
        'no_symbol': 19,             # NO_SYMBOL
        'flowchart_process': 61,     # FLOWCHART_PROCESS
        'flowchart_decision': 63,    # FLOWCHART_DECISION
        'flowchart_data': 64,        # FLOWCHART_DATA
        'flowchart_document': 67     # FLOWCHART_DOCUMENT
    }
    
    # Check if shape type is valid before trying to use it
    shape_type_lower = str(shape_type).lower()
    if shape_type_lower not in shape_type_map:
        available_shapes = ', '.join(sorted(shape_type_map.keys()))
        raise ValueError(f"Unsupported shape type: '{shape_type}'. Available shape types: {available_shapes}")
    
    # Get the integer value for the shape type
    shape_value = shape_type_map[shape_type_lower]
    
    # Create the shape using the direct integer value
    try:
        # The integer value is passed directly to add_shape
        shape = slide.shapes.add_shape(
            shape_value, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        return shape
    except Exception as e:
        raise ValueError(f"Failed to create '{shape_type}' shape using direct value {shape_value}: {str(e)}")

# ---- File download endpoint ----

@app.custom_route("/files/{filename}", methods=["GET"])
async def download_file(request: Request) -> FileResponse:
    """Serve saved presentation files for download."""
    filename = request.path_params["filename"]

    # Path traversal protection
    if "/" in filename or "\\" in filename or ".." in filename:
        return JSONResponse({"error": "Invalid filename"}, status_code=400)

    file_path = OUTPUT_DIR / filename
    resolved = file_path.resolve()

    # Ensure the resolved path is still within OUTPUT_DIR
    if not str(resolved).startswith(str(OUTPUT_DIR)):
        return JSONResponse({"error": "Invalid filename"}, status_code=400)

    if not resolved.is_file():
        return JSONResponse({"error": "File not found"}, status_code=404)

    return FileResponse(
        path=str(resolved),
        filename=filename,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ---- Custom presentation management wrapper ----

class PresentationManager:
    """Wrapper to handle presentation state updates."""
    
    def __init__(self, presentations_dict):
        self.presentations = presentations_dict
    
    def store_presentation(self, pres, pres_id):
        """Store a presentation and set it as current."""
        self.presentations[pres_id] = pres
        set_current_presentation_id(pres_id)
        return pres_id

# ---- Register Tools ----

# Create presentation manager wrapper
presentation_manager = PresentationManager(presentations)

# Wrapper functions to handle state management
def create_presentation_wrapper(original_func):
    """Wrapper to handle presentation creation with state management."""
    def wrapper(*args, **kwargs):
        result = original_func(*args, **kwargs)
        if "presentation_id" in result and result["presentation_id"] in presentations:
            set_current_presentation_id(result["presentation_id"])
        return result
    return wrapper

def open_presentation_wrapper(original_func):
    """Wrapper to handle presentation opening with state management."""
    def wrapper(*args, **kwargs):
        result = original_func(*args, **kwargs)
        if "presentation_id" in result and result["presentation_id"] in presentations:
            set_current_presentation_id(result["presentation_id"])
        return result
    return wrapper

# Register all tool modules
register_presentation_tools(
    app,
    presentations,
    get_current_presentation_id,
    get_template_search_directories,
    output_dir=OUTPUT_DIR,
    get_download_url_fn=get_download_url,
    transport_mode_fn=get_transport_mode,
)

register_content_tools(
    app,
    presentations,
    get_current_presentation_id,
    validate_parameters,
    is_positive,
    is_non_negative,
    is_in_range,
    is_valid_rgb
)

register_structural_tools(
    app,
    presentations,
    get_current_presentation_id,
    validate_parameters,
    is_positive,
    is_non_negative,
    is_in_range,
    is_valid_rgb,
    add_shape_direct
)

register_professional_tools(
    app,
    presentations,
    get_current_presentation_id
)

register_template_tools(
    app,
    presentations,
    get_current_presentation_id
)

register_hyperlink_tools(
    app,
    presentations,
    get_current_presentation_id,
    validate_parameters,
    is_positive,
    is_non_negative,
    is_in_range,
    is_valid_rgb
)

register_chart_tools(
    app,
    presentations,
    get_current_presentation_id,
    validate_parameters,
    is_positive,
    is_non_negative,
    is_in_range,
    is_valid_rgb
)


register_connector_tools(
    app,
    presentations,
    get_current_presentation_id,
    validate_parameters,
    is_positive,
    is_non_negative,
    is_in_range,
    is_valid_rgb
)

register_master_tools(
    app,
    presentations,
    get_current_presentation_id,
    validate_parameters,
    is_positive,
    is_non_negative,
    is_in_range,
    is_valid_rgb
)

register_transition_tools(
    app,
    presentations,
    get_current_presentation_id,
    validate_parameters,
    is_positive,
    is_non_negative,
    is_in_range,
    is_valid_rgb
)


# ---- Additional Utility Tools ----

@app.tool()
def list_presentations() -> Dict:
    """List all loaded presentations."""
    return {
        "presentations": [
            {
                "id": pres_id,
                "slide_count": len(pres.slides),
                "is_current": pres_id == current_presentation_id
            }
            for pres_id, pres in presentations.items()
        ],
        "current_presentation_id": current_presentation_id,
        "total_presentations": len(presentations)
    }

@app.tool()
def switch_presentation(presentation_id: str) -> Dict:
    """Switch to a different loaded presentation."""
    if presentation_id not in presentations:
        return {
            "error": f"Presentation '{presentation_id}' not found. Available presentations: {list(presentations.keys())}"
        }
    
    global current_presentation_id
    old_id = current_presentation_id
    current_presentation_id = presentation_id
    
    return {
        "message": f"Switched from presentation '{old_id}' to '{presentation_id}'",
        "previous_presentation_id": old_id,
        "current_presentation_id": current_presentation_id
    }

@app.tool()
def get_server_info() -> Dict:
    """Get information about the MCP server."""
    return {
        "name": "PowerPoint MCP Server - Enhanced Edition",
        "version": "2.1.0",
        "total_tools": 32,  # Organized into 11 specialized modules
        "loaded_presentations": len(presentations),
        "current_presentation": current_presentation_id,
        "features": [
            "Presentation Management (7 tools)",
            "Content Management (6 tools)", 
            "Template Operations (7 tools)",
            "Structural Elements (4 tools)",
            "Professional Design (3 tools)",
            "Specialized Features (5 tools)"
        ],
        "improvements": [
            "32 specialized tools organized into 11 focused modules",
            "68+ utility functions across 7 organized utility modules",
            "Enhanced parameter handling and validation",
            "Unified operation interfaces with comprehensive coverage",
            "Advanced template system with auto-generation capabilities",
            "Professional design tools with multiple effects and styling",
            "Specialized features including hyperlinks, connectors, slide masters",
            "Dynamic text sizing and intelligent wrapping",
            "Advanced visual effects and styling",
            "Content-aware optimization and validation",
            "Complete PowerPoint lifecycle management",
            "Modular architecture for better maintainability"
        ],
        "new_enhanced_features": [
            "Hyperlink Management - Add, update, remove, and list hyperlinks in text",
            "Advanced Chart Data Updates - Replace chart data with new categories and series",
            "Advanced Text Run Formatting - Apply formatting to specific text runs",
            "Shape Connectors - Add connector lines and arrows between points",
            "Slide Master Management - Access and manage slide masters and layouts",
            "Slide Transitions - Basic transition management (placeholder for future)"
        ]
    }

def main(transport: str = "stdio"):
    global _transport_mode
    _transport_mode = transport

    if transport in ("http", "sse"):
        port = int(os.environ.get("PORT", 8000))
        app.settings.host = "0.0.0.0"
        app.settings.port = port
        # Allow Railway's public/private domains through DNS rebinding protection
        public_domain = os.environ.get("RAILWAY_PUBLIC_DOMAIN", "")
        private_domain = os.environ.get("RAILWAY_PRIVATE_DOMAIN", "")
        allowed_hosts = ["127.0.0.1:*", "localhost:*", "[::1]:*"]
        if public_domain:
            allowed_hosts.append(public_domain)
        if private_domain:
            allowed_hosts.append(private_domain)
        app.settings.transport_security.allowed_hosts = allowed_hosts

        if transport == "http":
            api_key = get_api_key()
            if api_key:
                import uvicorn

                api_key_header = get_api_key_header_name()
                asgi_app = _build_streamable_http_app()
                asgi_app.add_middleware(
                    APIKeyMiddleware,
                    api_key=api_key,
                    header_name=api_key_header,
                    exempt_paths=["/health", "/healthz"],
                )
                print(f"API key auth enabled for streamable-http via header '{api_key_header}'")
                uvicorn.run(asgi_app, host=app.settings.host, port=app.settings.port)
            else:
                app.run(transport="streamable-http")
        else:
            api_key = get_api_key()
            if api_key:
                import uvicorn

                api_key_header = get_api_key_header_name()
                try:
                    asgi_app = _build_sse_app()
                except RuntimeError as exc:
                    raise RuntimeError(
                        "SSE transport with PPTX_MCP_API_KEY requires ASGI SSE app builder support in FastMCP."
                    ) from exc

                asgi_app.add_middleware(
                    APIKeyMiddleware,
                    api_key=api_key,
                    header_name=api_key_header,
                    exempt_paths=["/health", "/healthz"],
                )
                print(f"API key auth enabled for sse via header '{api_key_header}'")
                uvicorn.run(asgi_app, host=app.settings.host, port=app.settings.port)
            else:
                app.run(transport="sse")
    else:
        app.run(transport='stdio')
        
if __name__ == "__main__":
    # Parse command line arguments
    parser = argparse.ArgumentParser(description="MCP Server for PowerPoint manipulation using python-pptx")

    parser.add_argument(
        "-t",
        "--transport",
        type=str,
        default="stdio",
        choices=["stdio", "http", "sse"],
        help="Transport method for the MCP server (default: stdio)"
    )

    args = parser.parse_args()
    main(args.transport)
