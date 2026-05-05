"""
Content management utilities for PowerPoint MCP Server.
Functions for slides, text, images, tables, charts, and shapes.
"""
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from typing import Dict, List, Tuple, Optional, Any
import tempfile
import os
import base64


def add_slide(presentation: Presentation, layout_index: int = 1) -> Tuple:
    """
    Add a slide to the presentation.
    
    Args:
        presentation: The Presentation object
        layout_index: Index of the slide layout to use
        
    Returns:
        A tuple containing the slide and its layout
    """
    layout = presentation.slide_layouts[layout_index]
    slide = presentation.slides.add_slide(layout)
    return slide, layout


def get_slide_info(slide, slide_index: int) -> Dict:
    """
    Get information about a specific slide.
    
    Args:
        slide: The slide object
        slide_index: Index of the slide
        
    Returns:
        Dictionary containing slide information
    """
    try:
        from pptx.util import Inches, Emu

        placeholders = []
        for placeholder in slide.placeholders:
            placeholder_info = {
                "idx": placeholder.placeholder_format.idx,
                "type": str(placeholder.placeholder_format.type),
                "name": placeholder.name
            }
            placeholders.append(placeholder_info)
        
        shapes = []
        for i, shape in enumerate(slide.shapes):
            left_emu = shape.left
            top_emu = shape.top
            width_emu = shape.width
            height_emu = shape.height

            left_in = float(Inches(1).emu / left_emu) if left_emu else 0
            top_in = float(Inches(1).emu / top_emu) if top_emu else 0
            width_in = float(width_emu) / Inches(1).emu if width_emu else 0
            height_in = float(height_emu) / Inches(1).emu if height_emu else 0

            shape_info = {
                "index": i,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": left_emu,
                "top": top_emu,
                "width": width_emu,
                "height": height_emu,
                "left_inches": round(left_in, 2),
                "top_inches": round(top_in, 2),
                "width_inches": round(width_in, 2),
                "height_inches": round(height_in, 2),
                "has_text": bool(hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text.strip()),
            }
            shapes.append(shape_info)
        
        # Slide dimensions in inches (from presentation)
        try:
            pres = slide.part.package.presentation_part.presentation
            slide_width_emu = pres.slide_width
            slide_height_emu = pres.slide_height
        except Exception:
            slide_width_emu = None
            slide_height_emu = None
            
        slide_width_in = float(slide_width_emu) / Inches(1).emu if slide_width_emu else 0
        slide_height_in = float(slide_height_emu) / Inches(1).emu if slide_height_emu else 0

        return {
            "slide_index": slide_index,
            "layout_name": slide.slide_layout.name,
            "slide_width": slide_width_emu,
            "slide_height": slide_height_emu,
            "slide_width_inches": round(slide_width_in, 2),
            "slide_height_inches": round(slide_height_in, 2),
            "placeholder_count": len(placeholders),
            "placeholders": placeholders,
            "shape_count": len(shapes),
            "shapes": shapes
        }
    except Exception as e:
        raise Exception(f"Failed to get slide info: {str(e)}")


def set_title(slide, title: str) -> None:
    """
    Set the title of a slide.
    
    Args:
        slide: The slide object
        title: The title text
    """
    if slide.shapes.title:
        slide.shapes.title.text = title


def populate_placeholder(slide, placeholder_idx: int, text: str) -> None:
    """
    Populate a placeholder with text.
    
    Args:
        slide: The slide object
        placeholder_idx: The index of the placeholder
        text: The text to add
    """
    placeholder = slide.placeholders[placeholder_idx]
    placeholder.text = text


def add_bullet_points(placeholder, bullet_points: List[str]) -> None:
    """
    Add bullet points to a placeholder.
    
    Args:
        placeholder: The placeholder object
        bullet_points: List of bullet point texts
    """
    text_frame = placeholder.text_frame
    text_frame.clear()
    
    for i, point in enumerate(bullet_points):
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0


def add_textbox(slide, left: float, top: float, width: float, height: float, text: str,
                font_size: int = None, font_name: str = None, bold: bool = None,
                italic: bool = None, underline: bool = None, 
                color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                alignment: str = None, vertical_alignment: str = None, 
                auto_fit: bool = True) -> Any:
    """
    Add a textbox to a slide with formatting options.
    
    Args:
        slide: The slide object
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        text: Text content
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        underline: Whether text should be underlined
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical alignment ('top', 'middle', 'bottom')
        auto_fit: Whether to auto-fit text
        
    Returns:
        The created textbox shape
    """
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    
    textbox.text_frame.text = text
    
    # Apply formatting if provided
    if any([font_size, font_name, bold, italic, underline, color, bg_color, alignment, vertical_alignment]):
        format_text_advanced(
            textbox.text_frame,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            underline=underline,
            color=color,
            bg_color=bg_color,
            alignment=alignment,
            vertical_alignment=vertical_alignment
        )
    
    return textbox


def format_text(text_frame, font_size: int = None, font_name: str = None, 
                bold: bool = None, italic: bool = None, color: Tuple[int, int, int] = None,
                alignment: str = None) -> None:
    """
    Format text in a text frame.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
    """
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    for paragraph in text_frame.paragraphs:
        if alignment and alignment in alignment_map:
            paragraph.alignment = alignment_map[alignment]
            
        for run in paragraph.runs:
            font = run.font
            
            if font_size is not None:
                font.size = Pt(font_size)
            if font_name is not None:
                font.name = font_name
            if bold is not None:
                font.bold = bold
            if italic is not None:
                font.italic = italic
            if color is not None:
                r, g, b = color
                font.color.rgb = RGBColor(r, g, b)


def format_text_advanced(text_frame, font_size: int = None, font_name: str = None, 
                        bold: bool = None, italic: bool = None, underline: bool = None,
                        color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                        alignment: str = None, vertical_alignment: str = None) -> Dict:
    """
    Advanced text formatting with comprehensive options.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        underline: Whether text should be underlined
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical alignment ('top', 'middle', 'bottom')
    
    Returns:
        Dictionary with formatting results
    """
    result = {
        'success': True,
        'warnings': []
    }
    
    try:
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }

        vertical_alignment_map = {
            'top': MSO_VERTICAL_ANCHOR.TOP,
            'middle': MSO_VERTICAL_ANCHOR.MIDDLE,
            'bottom': MSO_VERTICAL_ANCHOR.BOTTOM
        }
        
        # Enable text wrapping
        text_frame.word_wrap = True

        if vertical_alignment and vertical_alignment in vertical_alignment_map:
            text_frame.vertical_anchor = vertical_alignment_map[vertical_alignment]
        
        # Apply formatting to all paragraphs and runs
        for paragraph in text_frame.paragraphs:
            if alignment and alignment in alignment_map:
                paragraph.alignment = alignment_map[alignment]
            
            for run in paragraph.runs:
                font = run.font
                
                if font_size is not None:
                    font.size = Pt(font_size)
                if font_name is not None:
                    font.name = font_name
                if bold is not None:
                    font.bold = bold
                if italic is not None:
                    font.italic = italic
                if underline is not None:
                    font.underline = underline
                if color is not None:
                    r, g, b = color
                    font.color.rgb = RGBColor(r, g, b)
                if bg_color is not None:
                    try:
                        r, g, b = bg_color
                        font.highlight_color.rgb = RGBColor(r, g, b)
                    except Exception:
                        pass
        
        return result
        
    except Exception as e:
        result['success'] = False
        result['error'] = str(e)
        return result


def add_image(slide, image_path: str, left: float, top: float, width: float = None, height: float = None) -> Any:
    """
    Add an image to a slide.
    
    Args:
        slide: The slide object
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional)
        height: Height in inches (optional)
        
    Returns:
        The created image shape
    """
    if width is not None and height is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), Inches(width), Inches(height)
        )
    elif width is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), Inches(width)
        )
    elif height is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), height=Inches(height)
        )
    else:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top)
        )


def add_table(slide, rows: int, cols: int, left: float, top: float, width: float, height: float) -> Any:
    """
    Add a table to a slide.
    
    Args:
        slide: The slide object
        rows: Number of rows
        cols: Number of columns
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        The created table shape
    """
    return slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )


def format_table_cell(cell, font_size: int = None, font_name: str = None, 
                     bold: bool = None, italic: bool = None, 
                     color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                     alignment: str = None, vertical_alignment: str = None) -> None:
    """
    Format a table cell.
    
    Args:
        cell: The table cell object
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment
        vertical_alignment: Vertical alignment
    """
    # Format text
    if any([font_size, font_name, bold, italic, color, alignment]):
        format_text_advanced(
            cell.text_frame,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            color=color,
            alignment=alignment
        )
    
    # Set background color
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*bg_color)


def add_chart(slide, chart_type: str, left: float, top: float, width: float, height: float,
              categories: List[str], series_names: List[str], series_values: List[List[float]]) -> Any:
    """
    Add a chart to a slide.
    
    Args:
        slide: The slide object
        chart_type: Type of chart ('column', 'bar', 'line', 'pie', etc.)
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: List of category names
        series_names: List of series names
        series_values: List of value lists for each series
        
    Returns:
        The created chart object
    """
    # Map chart type names to enum values
    chart_type_map = {
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
        'line': XL_CHART_TYPE.LINE,
        'line_markers': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'area': XL_CHART_TYPE.AREA,
        'stacked_area': XL_CHART_TYPE.AREA_STACKED,
        'scatter': XL_CHART_TYPE.XY_SCATTER,
        'radar': XL_CHART_TYPE.RADAR,
        'radar_markers': XL_CHART_TYPE.RADAR_MARKERS
    }
    
    xl_chart_type = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    for i, series_name in enumerate(series_names):
        if i < len(series_values):
            chart_data.add_series(series_name, series_values[i])
    
    # Add chart to slide
    chart_shape = slide.shapes.add_chart(
        xl_chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    )
    
    return chart_shape.chart


def format_chart(chart, has_legend: bool = True, legend_position: str = 'right',
                has_data_labels: bool = False, title: str = None,
                x_axis_title: str = None, y_axis_title: str = None,
                color_scheme: str = None) -> None:
    """
    Format a chart with various options.
    
    Args:
        chart: The chart object
        has_legend: Whether to show legend
        legend_position: Position of legend ('right', 'top', 'bottom', 'left')
        has_data_labels: Whether to show data labels
        title: Chart title
        x_axis_title: X-axis title
        y_axis_title: Y-axis title
        color_scheme: Color scheme to apply
    """
    try:
        # Set chart title
        if title:
            chart.chart_title.text_frame.text = title
        
        # Configure legend
        if has_legend:
            chart.has_legend = True
            # Note: Legend position setting may vary by chart type
        else:
            chart.has_legend = False
        
        # Configure data labels
        if has_data_labels:
            for series in chart.series:
                series.has_data_labels = True
        
        # Set axis titles if available
        try:
            if x_axis_title and hasattr(chart, 'category_axis'):
                chart.category_axis.axis_title.text_frame.text = x_axis_title
            if y_axis_title and hasattr(chart, 'value_axis'):
                chart.value_axis.axis_title.text_frame.text = y_axis_title
        except:
            pass  # Axis titles may not be available for all chart types
            
    except Exception:
        pass  # Graceful degradation for chart formatting


def _alignment_to_string(alignment) -> Optional[str]:
    if alignment is None:
        return None

    alignment_map = {
        PP_ALIGN.LEFT: 'left',
        PP_ALIGN.CENTER: 'center',
        PP_ALIGN.RIGHT: 'right',
        PP_ALIGN.JUSTIFY: 'justify'
    }

    return alignment_map.get(alignment, str(alignment))


def _vertical_anchor_to_string(anchor) -> Optional[str]:
    if anchor is None:
        return None

    anchor_map = {
        MSO_VERTICAL_ANCHOR.TOP: 'top',
        MSO_VERTICAL_ANCHOR.MIDDLE: 'middle',
        MSO_VERTICAL_ANCHOR.BOTTOM: 'bottom'
    }

    return anchor_map.get(anchor, str(anchor))


def _rgb_to_list(rgb) -> Optional[List[int]]:
    if rgb is None:
        return None

    try:
        return [int(rgb[0]), int(rgb[1]), int(rgb[2])]
    except Exception:
        return None


def _safe_get_highlight_color(font) -> Optional[List[int]]:
    try:
        highlight = font.highlight_color
    except Exception:
        return None
    if highlight is None:
        return None

    try:
        return _rgb_to_list(highlight.rgb)
    except Exception:
        return None


def _is_placeholder(shape) -> bool:
    try:
        _ = shape.placeholder_format
        return True
    except Exception:
        return False


def _extract_text_frame_formatting(text_frame) -> Dict:
    formatting = {
        "word_wrap": getattr(text_frame, 'word_wrap', None),
        "vertical_alignment": _vertical_anchor_to_string(getattr(text_frame, 'vertical_anchor', None)),
        "paragraphs": []
    }

    try:
        paragraphs = list(text_frame.paragraphs)
    except Exception:
        return formatting

    for paragraph_index, paragraph in enumerate(paragraphs):
        paragraph_info = {
            "index": paragraph_index,
            "alignment": _alignment_to_string(getattr(paragraph, 'alignment', None)),
            "runs": []
        }

        try:
            runs = list(paragraph.runs)
        except Exception:
            runs = []

        for run_index, run in enumerate(runs):
            font = getattr(run, 'font', None)
            font_size = None
            if font is not None and getattr(font, 'size', None) is not None:
                try:
                    font_size = int(font.size.pt)
                except Exception:
                    font_size = None

            color = None
            if font is not None and getattr(font, 'color', None) is not None:
                try:
                    color = _rgb_to_list(font.color.rgb)
                except Exception:
                    color = None

            bg_color = _safe_get_highlight_color(font) if font is not None else None

            paragraph_info["runs"].append({
                "index": run_index,
                "text": getattr(run, 'text', ''),
                "font_name": font.name if font is not None else None,
                "font_size": font_size,
                "bold": font.bold if font is not None else None,
                "italic": font.italic if font is not None else None,
                "underline": font.underline if font is not None else None,
                "color": color,
                "bg_color": bg_color
            })

        formatting["paragraphs"].append(paragraph_info)

    return formatting


def extract_slide_text_content(slide) -> Dict:
    """
    Extract all text content from a slide including placeholders and text shapes.
    
    Args:
        slide: The slide object to extract text from
        
    Returns:
        Dictionary containing all text content organized by source type
    """
    try:
        text_content = {
            "slide_title": "",
            "placeholders": [],
            "text_shapes": [],
            "table_text": [],
            "table_cells": [],
            "all_text_combined": ""
        }
        
        all_texts = []
        
        # Extract title from slide if available
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            try:
                title_text = slide.shapes.title.text_frame.text.strip()
                if title_text:
                    text_content["slide_title"] = title_text
                    all_texts.append(title_text)
            except:
                pass
        
        # Extract text from all shapes
        for i, shape in enumerate(slide.shapes):
            shape_text_info = {
                "shape_index": i,
                "shape_name": shape.name,
                "shape_type": str(shape.shape_type),
                "text": ""
            }
            
            try:
                # Check if shape has text frame
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        shape_text_info["formatting"] = _extract_text_frame_formatting(shape.text_frame)
                        shape_text_info["text"] = text
                        all_texts.append(text)
                        
                        # Categorize by shape type
                        if _is_placeholder(shape):
                            # This is a placeholder
                            placeholder_info = shape_text_info.copy()
                            placeholder_info["placeholder_type"] = str(shape.placeholder_format.type)
                            placeholder_info["placeholder_idx"] = shape.placeholder_format.idx
                            text_content["placeholders"].append(placeholder_info)
                        else:
                            # This is a regular text shape
                            text_content["text_shapes"].append(shape_text_info)
                
                # Extract text from tables
                elif hasattr(shape, 'table'):
                    table_texts = []
                    table_cell_details = []
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        row_texts = []
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text_frame.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                                all_texts.append(cell_text)
                                table_cell_details.append({
                                    "row": row_idx,
                                    "col": col_idx,
                                    "text": cell_text,
                                    "formatting": _extract_text_frame_formatting(cell.text_frame)
                                })
                        if row_texts:
                            table_texts.append({
                                "row": row_idx,
                                "cells": row_texts
                            })
                    
                    if table_texts:
                        text_content["table_text"].append({
                            "shape_index": i,
                            "shape_name": shape.name,
                            "table_content": table_texts
                        })

                    if table_cell_details:
                        text_content["table_cells"].append({
                            "shape_index": i,
                            "shape_name": shape.name,
                            "cells": table_cell_details
                        })
                        
            except Exception as e:
                # Skip shapes that can't be processed
                continue
        
        # Combine all text
        text_content["all_text_combined"] = "\n".join(all_texts)
        
        return {
            "success": True,
            "text_content": text_content,
            "total_text_shapes": len(text_content["placeholders"]) + len(text_content["text_shapes"]),
            "has_title": bool(text_content["slide_title"]),
            "has_tables": len(text_content["table_text"]) > 0
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": f"Failed to extract text content: {str(e)}",
            "text_content": None
        }